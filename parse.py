import PyPDF2
import pandas as pd
from docx import Document
import win32com.client as win32
import os
import json
from zipfile import ZipFile
import xml.etree.ElementTree as ET
from pptx import Presentation
import tempfile

class FileParser:
    def __init__(self, root_dir):
        self.root_dir = root_dir

    def read_txt_file(self,file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
        
    def read_pdf_file(self, file_path):
        with open(file_path, 'rb') as file:
            render = PyPDF2.PdfReader(file)
            text=''
            for page in render.pages:
                text += page.extract_text()
            return text
        
    def read_doc_file(self, file_path):
        word = win32.Dispatch("Word.Application")
        doc = word.Documents.Open(file_path)
        text = doc.Range().Text
        doc.Close()
        word.Quit()
        return text

    def read_docx_file(self, file_path):
        doc = Document(file_path)
        text = '\n'.join([para.text for para in doc.paragraphs])
        return text

    def read_csv_file(self, file_path):
        return pd.read_csv(file_path).to_string()

    def read_xlsx_file(self, file_path):
        return pd.read_excel(file_path)

    def read_json_file(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            return json.dumps(json.load(file), indent=4)
        
    def read_xml_file(self, file_path):
        tree = ET.parse(file_path)
        root = tree.getroot()
        return ET.tostring(root, encoding='unicode')

    def read_pptx_file(self, file_path):
        prs = Presentation(file_path)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text

    def read_zip_file(self, file_path):
        extracted_data = {}
        with tempfile.TemporaryDirectory() as temp_dir:
            with ZipFile(file_path, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)
                extracted_data.update(self.parse_files(temp_dir))
        return extracted_data

    def parse_files(self, root_dir=None):
        if root_dir is None:
            root_dir = self.root_dir
        file_data = {}
        for root, _, files in os.walk(root_dir):
            for file in files:
                file_path = os.path.join(root, file)
                extension = os.path.splitext(file)[1].lower()
                try:            
                    if extension == '.txt':
                        file_data[file_path] = self.read_txt_file(file_path)
                    elif extension == '.pdf':
                        file_data[file_path] = self.read_pdf_file(file_path)
                    elif extension == '.doc':
                        file_data[file_path] = self.read_doc_file(file_path)
                    elif extension == '.docx':
                        file_data[file_path] = self.read_docx_file(file_path)
                    elif extension == '.csv':
                        file_data[file_path] = self.read_csv_file(file_path)
                    elif extension == '.xlsx':
                        file_data[file_path] = self.read_xlsx_file(file_path)
                    elif extension == '.json':
                        file_data[file_path] = self.read_json_file(file_path)
                    elif extension == '.xml':
                        file_data[file_path] = self.read_xml_file(file_path)
                    elif extension == '.pptx':
                        file_data[file_path] = self.read_pptx_file(file_path)
                    elif extension == '.zip':
                        file_data[file_path] = self.read_zip_file(file_path)
                except Exception as e:
                    print(f"Error reading {file_path}: {e}")
        return file_data

parser = FileParser(r"C:\Users\swast\OneDrive\Desktop\files")
parsed_files = parser.parse_files()
print(parsed_files)

# root_dir = r"C:\Users\swast\OneDrive\Desktop\files"
# parsed_files = parse_files(root_dir)
# for file_path, content in parsed_files.items():
#     print(f"File: {file_path}")
#     if isinstance(content, dict):
#         print(f"Content: {content}")
#     else:
#         print(f"Content: {content[:200]}...")                  

